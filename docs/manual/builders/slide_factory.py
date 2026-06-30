"""
OmniMob Diamond - Slide Factory.
Funcoes que produzem os tipos de slide do padrao Diamond.
Espelha a estrutura do Padrao Gold da Enermac com identidade OmniMob.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from . import theme as t


def _blank_slide(prs):
    layout = prs.slide_layouts[6]
    return prs.slides.add_slide(layout)


def _set_fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def _no_line(shape):
    shape.line.fill.background()


def _solid_line(shape, rgb, width_pt=0.75):
    shape.line.color.rgb = rgb
    shape.line.width = Pt(width_pt)


def _add_rect(slide, x, y, w, h, fill, line=None, line_w=0):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    _set_fill(rect, fill)
    if line is None:
        _no_line(rect)
    else:
        _solid_line(rect, line, line_w)
    rect.text_frame.text = ""
    return rect


def _add_round_rect(slide, x, y, w, h, fill, line=None, line_w=0.75):
    r = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    _set_fill(r, fill)
    if line is None:
        _no_line(r)
    else:
        _solid_line(r, line, line_w)
    return r


def _add_text(slide, x, y, w, h, text, *,
              size=None, color=None, bold=False,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
              font=None, italic=False):
    if size is None:
        size = t.SIZE_BODY
    if color is None:
        color = t.TEXT_WHITE
    if font is None:
        font = t.FONT_FAMILY

    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor

    if isinstance(text, str):
        lines = text.split("\n")
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            r = p.add_run()
            r.text = line
            r.font.name = font
            r.font.size = size
            r.font.bold = bold
            r.font.italic = italic
            r.font.color.rgb = color
    elif isinstance(text, list):
        for i, para in enumerate(text):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            if isinstance(para, str):
                r = p.add_run()
                r.text = para
                r.font.name = font
                r.font.size = size
                r.font.bold = bold
                r.font.italic = italic
                r.font.color.rgb = color
            else:
                for run_spec in para:
                    if isinstance(run_spec, str):
                        r = p.add_run()
                        r.text = run_spec
                        r.font.name = font
                        r.font.size = size
                        r.font.bold = bold
                        r.font.color.rgb = color
                    else:
                        txt, opts = run_spec
                        r = p.add_run()
                        r.text = txt
                        r.font.name = opts.get("font", font)
                        r.font.size = opts.get("size", size)
                        r.font.bold = opts.get("bold", bold)
                        r.font.italic = opts.get("italic", italic)
                        r.font.color.rgb = opts.get("color", color)
    return tb


def _draw_logo(slide, x, y, w, h, *, on_dark=True):
    dot_size = h * 0.55
    dot_y = y + (h - dot_size) / 2
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, dot_y, dot_size, dot_size)
    _set_fill(dot, t.BRAND_GREEN)
    _no_line(dot)

    txt_x = x + dot_size + Inches(0.08)
    txt_w = w - dot_size - Inches(0.08)
    color = t.TEXT_WHITE if on_dark else t.BG_DARK
    _add_text(slide, txt_x, y, txt_w, h,
              "OmniMob", size=Pt(20), bold=True,
              color=color, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.LEFT)


def cover_slide(prs, *, module_num, category, title, subtitle, audience):
    slide = _blank_slide(prs)
    _add_rect(slide, 0, 0, t.SLIDE_W, t.SLIDE_H, t.BG_DARK)
    _add_rect(slide, 0, 0, t.COVER_STRIPE_W, t.SLIDE_H, t.BRAND_GREEN)

    band_y = Inches(2.4)
    band_h = Inches(2.8)
    _add_rect(slide, 0, band_y, t.SLIDE_W, band_h, t.BG_CARD)
    _add_rect(slide, 0, band_y, t.SLIDE_W, Inches(0.045), t.ACCENT_GOLD)

    _draw_logo(slide, Inches(0.7), Inches(0.55), Inches(2.2), Inches(0.55))

    _add_text(slide, Inches(0.7), Inches(1.18), Inches(6), Inches(0.3),
              "MANUAL OFICIAL — PADRAO DIAMOND",
              size=Pt(11), bold=True, color=t.ACCENT_GOLD)

    _add_text(slide, Inches(0.7), band_y + Inches(0.30), Inches(11), Inches(0.4),
              "Modulo " + ("%02d" % module_num) + " — " + category,
              size=Pt(16), bold=True, color=t.BRAND_GREEN)

    _add_text(slide, Inches(0.7), band_y + Inches(0.70), Inches(12), Inches(1.1),
              title, size=t.SIZE_DISPLAY, bold=True, color=t.TEXT_WHITE)

    _add_text(slide, Inches(0.7), band_y + Inches(1.85), Inches(12), Inches(0.5),
              subtitle, size=Pt(15), color=t.TEXT_MUTED)

    aud_label, aud_color = audience
    _add_text(slide, Inches(0.7), Inches(5.6), Inches(6), Inches(0.4),
              aud_label, size=Pt(15), bold=True, color=aud_color)

    _add_rect(slide, Inches(0.7), Inches(6.0), Inches(0.8), Inches(0.04), t.ACCENT_GOLD)

    _add_text(slide, Inches(0.7), Inches(7.05), Inches(12), Inches(0.3),
              t.FOOTER_TEXT + "  —  " + t.BRAND_TAGLINE,
              size=t.SIZE_CAPTION, color=t.TEXT_DIM)

    return slide


def _header(slide, title, subtitle=None, audience=None, module_num=None):
    _add_rect(slide, 0, 0, t.SLIDE_W, t.HEADER_H, t.BG_DARK)
    _add_rect(slide, 0, t.HEADER_H, t.SLIDE_W, t.HEADER_STRIPE_H, t.BRAND_GREEN)

    _add_text(slide, Inches(0.55), Inches(0.18), Inches(10), Inches(0.55),
              title, size=Pt(24), bold=True, color=t.TEXT_WHITE)

    if subtitle:
        _add_text(slide, Inches(0.55), Inches(0.70), Inches(10), Inches(0.32),
                  subtitle, size=Pt(12), color=t.TEXT_MUTED)

    if module_num is not None and audience:
        aud_label, aud_color = audience
        _add_text(slide,
                  t.SLIDE_W - Inches(3.7), Inches(0.22), Inches(3.5), Inches(0.32),
                  "MODULO " + ("%02d" % module_num) + "   •   " + aud_label,
                  size=Pt(10), bold=True, color=aud_color, align=PP_ALIGN.RIGHT)


def _footer(slide, page_num=None, total=None):
    _add_rect(slide, Inches(0.55), t.SLIDE_H - Inches(0.40),
              t.SLIDE_W - Inches(1.1), Emu(9525), t.BORDER)

    _add_text(slide, Inches(0.55), t.SLIDE_H - Inches(0.32),
              Inches(8), Inches(0.25),
              t.FOOTER_TEXT, size=t.SIZE_TINY, color=t.TEXT_DIM)

    if page_num is not None:
        page_txt = str(page_num) + (" / " + str(total) if total else "")
        _add_text(slide,
                  t.SLIDE_W - Inches(2), t.SLIDE_H - Inches(0.32),
                  Inches(1.5), Inches(0.25),
                  page_txt, size=t.SIZE_TINY, color=t.TEXT_DIM,
                  align=PP_ALIGN.RIGHT)


def hook_slide(prs, *, module_num, audience, headline, hook_phrase, objectives):
    slide = _blank_slide(prs)
    _add_rect(slide, 0, 0, t.SLIDE_W, t.SLIDE_H, t.BG_DARK)
    _header(slide, "Por que este modulo importa",
            subtitle=headline, audience=audience, module_num=module_num)

    box_y = Inches(1.50)
    box_h = Inches(1.40)
    _add_round_rect(slide, Inches(0.55), box_y,
                    t.SLIDE_W - Inches(1.1), box_h,
                    t.BG_CARD, line=t.ACCENT_GOLD, line_w=2.0)

    _add_text(slide, Inches(0.9), box_y, t.SLIDE_W - Inches(1.8), box_h,
              hook_phrase,
              size=Pt(18), bold=True, color=t.TEXT_WHITE,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    list_y = Inches(3.30)
    _add_text(slide, Inches(0.55), list_y, Inches(11), Inches(0.4),
              "Ao final deste modulo, voce sera capaz de:",
              size=Pt(15), bold=True, color=t.TEXT_WHITE)

    y = list_y + Inches(0.55)
    for obj in objectives[:5]:
        _add_text(slide, Inches(0.65), y, Inches(0.4), Inches(0.40),
                  t.ICON_CHECK, size=Pt(16), bold=True, color=t.BRAND_GREEN)
        _add_text(slide, Inches(1.05), y, Inches(11.6), Inches(0.40),
                  obj, size=t.SIZE_BODY_SM, color=t.TEXT_WHITE,
                  anchor=MSO_ANCHOR.TOP)
        y += Inches(0.50)

    _footer(slide, page_num=2)
    return slide


def content_slide(prs, *, module_num, audience, title, subtitle, body,
                  page_num=None):
    slide = _blank_slide(prs)
    _add_rect(slide, 0, 0, t.SLIDE_W, t.SLIDE_H, t.BG_DARK)
    _header(slide, title, subtitle=subtitle, audience=audience, module_num=module_num)

    card_y = t.HEADER_H + t.HEADER_STRIPE_H + Inches(0.30)
    card_h = t.SLIDE_H - card_y - Inches(0.65)
    _add_round_rect(slide,
                    Inches(0.55), card_y,
                    t.SLIDE_W - Inches(1.1), card_h,
                    t.BG_CARD, line=t.BORDER, line_w=0.75)

    y = card_y + Inches(0.35)
    inner_w = t.SLIDE_W - Inches(1.7)
    max_items = 8
    for item in body[:max_items]:
        if isinstance(item, str):
            _add_text(slide, Inches(0.85), y, Inches(0.30), Inches(0.40),
                      t.ICON_DOT, size=Pt(14), color=t.BRAND_GREEN, bold=True)
            _add_text(slide, Inches(1.15), y, inner_w - Inches(0.30), Inches(0.40),
                      item, size=t.SIZE_BODY_SM, color=t.TEXT_WHITE)
            y += Inches(0.48)
        elif isinstance(item, dict):
            label = item.get("label", "")
            desc = item.get("desc", "")
            _add_text(slide, Inches(0.85), y, Inches(0.30), Inches(0.40),
                      t.ICON_ARROW, size=Pt(14), color=t.ACCENT_GOLD, bold=True)
            runs = [[
                (label, {"bold": True, "color": t.BRAND_GREEN, "size": t.SIZE_BODY_SM}),
                ("  —  ", {"color": t.TEXT_MUTED, "size": t.SIZE_BODY_SM}),
                (desc, {"color": t.TEXT_WHITE, "size": t.SIZE_BODY_SM})
            ]]
            _add_text(slide, Inches(1.15), y, inner_w - Inches(0.30), Inches(0.40),
                      runs, size=t.SIZE_BODY_SM, color=t.TEXT_WHITE)
            y += Inches(0.55)

    _footer(slide, page_num=page_num)
    return slide


def screenshot_slide(prs, *, module_num, audience, title, subtitle,
                     image_path=None, caption=None, hotspot_notes=None,
                     page_num=None):
    slide = _blank_slide(prs)
    _add_rect(slide, 0, 0, t.SLIDE_W, t.SLIDE_H, t.BG_DARK)
    _header(slide, title, subtitle=subtitle, audience=audience, module_num=module_num)

    img_x = Inches(0.55)
    img_y = t.HEADER_H + t.HEADER_STRIPE_H + Inches(0.30)
    img_w = Inches(8.4)
    img_h = Inches(4.8)

    _add_round_rect(slide, img_x, img_y, img_w, img_h,
                    t.BG_CARD, line=t.BRAND_GREEN, line_w=1.5)

    if image_path:
        try:
            slide.shapes.add_picture(str(image_path),
                                     img_x + Inches(0.10),
                                     img_y + Inches(0.10),
                                     width=img_w - Inches(0.20),
                                     height=img_h - Inches(0.20))
        except Exception as e:
            _add_text(slide, img_x, img_y, img_w, img_h,
                      "[Print nao carregado: " + str(e) + "]",
                      size=t.SIZE_BODY, color=t.ERROR,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    else:
        _add_text(slide, img_x, img_y, img_w, img_h,
                  "[ PRINT REAL DO SISTEMA ]\nA ser capturado pelo pipeline Playwright",
                  size=Pt(14), color=t.TEXT_MUTED,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    if caption:
        _add_text(slide, img_x, img_y + img_h + Inches(0.08),
                  img_w, Inches(0.35),
                  caption, size=t.SIZE_CAPTION, italic=True,
                  color=t.TEXT_MUTED, align=PP_ALIGN.CENTER)

    notes_x = Inches(9.20)
    notes_w = t.SLIDE_W - notes_x - Inches(0.55)
    notes_h = Inches(5.4)
    _add_round_rect(slide, notes_x, img_y, notes_w, notes_h,
                    t.BG_CARD, line=t.BORDER, line_w=0.75)

    _add_text(slide, notes_x + Inches(0.20), img_y + Inches(0.15),
              notes_w - Inches(0.40), Inches(0.35),
              "Pontos marcados", size=t.SIZE_CAPTION, bold=True,
              color=t.ACCENT_GOLD)

    if hotspot_notes:
        y = img_y + Inches(0.55)
        for i, note in enumerate(hotspot_notes[:8], 1):
            num_x = notes_x + Inches(0.20)
            num_size = Inches(0.30)
            bullet = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                            num_x, y, num_size, num_size)
            _set_fill(bullet, t.BRAND_GREEN)
            _no_line(bullet)
            tb = bullet.text_frame
            tb.margin_left = Emu(0)
            tb.margin_right = Emu(0)
            tb.margin_top = Emu(0)
            tb.margin_bottom = Emu(0)
            p = tb.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = str(i)
            r.font.size = Pt(10)
            r.font.bold = True
            r.font.color.rgb = t.BG_DARK
            r.font.name = t.FONT_FAMILY

            _add_text(slide, num_x + num_size + Inches(0.10), y - Inches(0.02),
                      notes_w - num_size - Inches(0.30), Inches(0.50),
                      note, size=Pt(10), color=t.TEXT_WHITE)
            y += Inches(0.50)

    _footer(slide, page_num=page_num)
    return slide


def recap_slide(prs, *, module_num, audience, title, takeaways, next_module=None,
                page_num=None):
    slide = _blank_slide(prs)
    _add_rect(slide, 0, 0, t.SLIDE_W, t.SLIDE_H, t.BG_DARK)
    _header(slide, "Recap — o que voce leva deste modulo",
            subtitle=title, audience=audience, module_num=module_num)

    y = Inches(1.5)
    n = min(len(takeaways), 4)
    if n == 0:
        n = 1
    card_w = (t.SLIDE_W - Inches(1.1) - Inches(0.3) * (n - 1)) / n
    card_h = Inches(2.8)
    x = Inches(0.55)
    for i, take in enumerate(takeaways[:n]):
        _add_round_rect(slide, x, y, card_w, card_h, t.BG_CARD,
                        line=t.BRAND_GREEN if i % 2 == 0 else t.ACCENT_GOLD,
                        line_w=1.2)
        _add_text(slide, x + Inches(0.20), y + Inches(0.15),
                  card_w, Inches(0.6),
                  "#" + ("%02d" % (i + 1)), size=Pt(36), bold=True,
                  color=(t.BRAND_GREEN if i % 2 == 0 else t.ACCENT_GOLD))
        _add_text(slide, x + Inches(0.20), y + Inches(0.90),
                  card_w - Inches(0.40), card_h - Inches(1.0),
                  take, size=Pt(13), color=t.TEXT_WHITE)
        x += card_w + Inches(0.30)

    if next_module:
        nx_y = Inches(4.7)
        _add_round_rect(slide, Inches(0.55), nx_y,
                        t.SLIDE_W - Inches(1.1), Inches(1.3),
                        t.BG_CARD, line=t.ACCENT_GOLD, line_w=1.5)
        _add_text(slide, Inches(0.85), nx_y + Inches(0.15),
                  Inches(11), Inches(0.4),
                  "Proximo passo", size=Pt(12), bold=True, color=t.ACCENT_GOLD)
        _add_text(slide, Inches(0.85), nx_y + Inches(0.50),
                  Inches(11), Inches(0.6),
                  next_module, size=Pt(16), bold=True, color=t.TEXT_WHITE)

    _footer(slide, page_num=page_num)
    return slide


def new_presentation():
    prs = Presentation()
    prs.slide_width = t.SLIDE_W
    prs.slide_height = t.SLIDE_H
    return prs